#!/usr/bin/env python3
"""
Update The Invisible Users presentation with new slides and content.

This script demonstrates how to programmatically modify PowerPoint presentations using
python-pptx library. It was created to add three specific changes to the presentation:
1. Update slide 4 to add a fifth agent type (Cowork)
2. Add a new slide about the Ally McBeal case after slide 9
3. Add a JSON-LD pricing example slide after slide 16

CUSTOMIZATION GUIDE FOR AI ASSISTANTS:
---------------------------------------
To adapt this script for other presentation updates:

1. UPDATE FILE PATHS (in main() function):
   - input_file: Path to source PPTX file
   - output_file: Path where updated PPTX will be saved

2. MODIFY SLIDE POSITIONS:
   - Slide indices are 0-based (slide 1 = index 0, slide 2 = index 1, etc.)
   - Change slide index numbers in function calls and duplicate_slide() calls

3. CUSTOMIZE CONTENT:
   - Update text strings in each function (titles, subtitles, bullet points)
   - Modify shape identification logic if template layout differs

4. ADD NEW FUNCTIONS:
   - Follow the pattern of existing functions (identify shapes, clear, add new content)
   - Call new functions from main() in the desired order

5. CHANGE FORMATTING:
   - Font: Modify font.name (e.g., "Arial", "Courier New", "Calibri")
   - Size: Modify font.size using Pt() (e.g., Pt(14), Pt(18))
   - Color: Use RGBColor(r, g, b) with values 0-255
   - Alignment: Use PP_ALIGN.LEFT, PP_ALIGN.CENTER, or PP_ALIGN.RIGHT

COMMON PATTERNS:
----------------
- Identify shape by content: Search for unique text in shape.text_frame.text
- Clear and replace: shape.text_frame.clear() then add new paragraphs
- Preserve formatting: Copy font properties from previous paragraph
- Duplicate slides: Use duplicate_slide(prs, index) from rearrange.py
- Move slides: Use prs.slides._sldIdLst to reorder slide elements

DEPENDENCIES:
-------------
- python-pptx: pip install python-pptx
- Helper functions from Claude Code PPTX skill (rearrange.py)

USAGE:
------
python scripts/update_invisible_users_presentation.py

The script will:
- Load the original presentation
- Make all modifications
- Save to a new file (preserves original)
- Print progress and final slide count

METADATA:
---------
Created: 2026-01-21
Purpose: Batch update presentation slides programmatically
Original presentation: The Invisible Users Presentation-final.pptx (23 slides)
Updated presentation: The Invisible Users Presentation-final-updated.pptx (25 slides)
Changes: +2 new slides, 1 modified slide
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from copy import deepcopy
import six

# Add path for helper functions from Claude Code PPTX skill
# This provides the duplicate_slide() function for cloning slides
sys.path.insert(0, str(Path.home() / '.claude/plugins/cache/anthropic-agent-skills/document-skills/69c0b1a06741/skills/pptx/scripts'))
from rearrange import duplicate_slide

def main():
    """
    Main entry point - orchestrates all presentation updates.

    CUSTOMIZATION: Update file paths here to work with different presentations.
    """
    input_file = "docs/scrapboard/The Invisible Users Presentation-final.pptx"
    output_file = "docs/scrapboard/The Invisible Users Presentation-final-updated.pptx"

    print(f"Loading presentation from: {input_file}")
    prs = Presentation(input_file)
    print(f"Original slide count: {len(prs.slides)}")

    # Change 1: Update Slide 4 (index 3) - Add Cowork as fifth agent type
    print("\n=== Change 1: Updating Slide 4 (Four Types → Five Types) ===")
    slide_4 = prs.slides[3]
    update_slide_4_add_cowork(slide_4)

    # Change 2: Add Ally McBeal slide after slide 9 (index 8)
    print("\n=== Change 2: Adding Ally McBeal slide after slide 9 ===")
    add_ally_mcbeal_slide(prs)
    print(f"After adding Ally McBeal slide: {len(prs.slides)} slides")

    # Change 3: Add JSON-LD pricing example slide after slide 16 (now index 16 after adding slide 10)
    print("\n=== Change 3: Adding JSON-LD example slide after slide 16 ===")
    add_json_ld_slide(prs)
    print(f"After adding JSON-LD slide: {len(prs.slides)} slides")

    # Save
    print(f"\nSaving updated presentation to: {output_file}")
    prs.save(output_file)
    print(f"✓ Successfully saved! Final slide count: {len(prs.slides)}")

def update_slide_4_add_cowork(slide):
    """
    Update slide 4 to add Cowork as fifth agent type.

    This function demonstrates how to:
    - Find and update title text via string replacement
    - Add a new bullet point to an existing list
    - Preserve formatting from previous paragraphs

    CUSTOMIZATION:
    - Change "Four Types" and "Five Types" strings to match your slide
    - Update identification strings ("Server-Side", "Local/On-Device")
    - Modify new bullet text and description
    - Adjust paragraph level (0=top level, 1=indented, etc.)

    Args:
        slide: pptx.slide.Slide object to modify
    """
    print("  Finding title and body shapes...")

    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue

        text = shape.text_frame.text

        # Update title if it contains "Four Types"
        if "Four Types" in text:
            print(f"  Found title shape: '{text}'")
            print("  Updating to 'Five Types of AI Agents'")
            shape.text_frame.text = text.replace("Four Types", "Five Types")

        # Add Cowork to body if it contains the list of agent types
        elif "Server-Side" in text and "Local/On-Device" in text:
            print("  Found body shape with agent types")
            print("  Adding: Agentic OS (Cowork)")

            # Add new paragraph for Cowork
            p = shape.text_frame.add_paragraph()
            p.text = "Agentic OS (Cowork) - Orchestrates multiple agents in unified environment"
            p.level = 0

            # Copy formatting from previous paragraph
            if len(shape.text_frame.paragraphs) > 1:
                prev_para = shape.text_frame.paragraphs[-2]
                if prev_para.runs:
                    prev_run = prev_para.runs[0]
                    if p.runs:
                        p.runs[0].font.size = prev_run.font.size
                        p.runs[0].font.name = prev_run.font.name
                        p.runs[0].font.color.rgb = prev_run.font.color.rgb

    print("  ✓ Slide 4 updated successfully")

def add_ally_mcbeal_slide(prs):
    """
    Add Ally McBeal case study slide after slide 9.

    This function demonstrates how to:
    - Duplicate an existing slide as a template
    - Move the duplicated slide to a specific position
    - Identify shapes by their content
    - Replace text in multiple shapes (title, subtitle, body)
    - Add multiple bullet points programmatically

    CUSTOMIZATION:
    - Change source slide index (currently 8 = slide 9)
    - Update target position (currently index 9 = position 10)
    - Modify identification strings to match your template
    - Update title, subtitle, and bullet point content
    - Adjust bullet point count as needed

    Args:
        prs: pptx.presentation.Presentation object
    """
    print("  Duplicating slide 9 (Toast Notifications)...")

    # Duplicate slide 9 (index 8)
    new_slide = duplicate_slide(prs, 8)
    print(f"  Created new slide at index {len(prs.slides) - 1}")

    # Move the duplicated slide to the correct position
    # New slides are added at the end, so we need to move it to position 10
    print("  Moving to position 10...")
    slides = prs.slides._sldIdLst  # Get internal slide list
    slide_element = slides[len(slides) - 1]  # Get the last slide (our new one)
    slides.remove(slide_element)  # Remove from end
    slides.insert(9, slide_element)  # Insert at index 9 (position 10)

    # Update content in the moved slide
    print("  Updating slide content...")
    slide = prs.slides[9]  # Now access it at its new position

    # Strategy: Identify shapes by their content, then replace text
    # This is more reliable than position-based identification
    title_shape = None
    subtitle_shape = None
    body_shape = None

    # First pass: identify which shape is which
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue  # Skip non-text shapes (images, etc.)

        text = shape.text_frame.text

        # Identify shapes by unique text from the template slide
        # CUSTOMIZATION: Change these strings to match your template
        if "Mistake #1" in text:
            title_shape = shape
        elif "Pattern That Keeps" in text:
            subtitle_shape = shape
        elif "Toast notifications vanish" in text or "showToast" in text:
            body_shape = shape

    # Second pass: update the identified shapes
    # Update title
    if title_shape:
        print("    Updating title")
        title_shape.text_frame.clear()  # Remove all existing text
        p = title_shape.text_frame.paragraphs[0]  # Get first paragraph
        p.text = "Real-World Consequences: The Ally McBeal Case"
        # CUSTOMIZATION: Change title text here

    # Update subtitle
    if subtitle_shape:
        print("    Updating subtitle")
        subtitle_shape.text_frame.clear()
        p = subtitle_shape.text_frame.paragraphs[0]
        p.text = "When metadata fails, AI agents hallucinate dangerously"
        # CUSTOMIZATION: Change subtitle text here

    # Update body with bullets
    if body_shape:
        print("    Updating body bullets")
        body_shape.text_frame.clear()

        # CUSTOMIZATION: Update bullet points here
        bullets = [
            "Lawyers caught citing fictional cases in court",
            "AI agents confused Ally McBeal TV scripts with legal precedents",
            "Without proper microdata/metadata distinguishing entertainment from legal docs",
            "Agents fabricate details that seem plausible but are dangerously incorrect"
        ]

        # Add each bullet as a paragraph
        for i, bullet_text in enumerate(bullets):
            if i == 0:
                p = body_shape.text_frame.paragraphs[0]  # Use existing first paragraph
            else:
                p = body_shape.text_frame.add_paragraph()  # Add new paragraphs
            p.text = bullet_text
            p.level = 0  # Set indent level (0=no indent, 1=first level indent, etc.)

    print("  ✓ Ally McBeal slide added successfully")

def add_json_ld_slide(prs):
    """
    Add JSON-LD example slide after slide 16 (now index 16 after previous addition).

    This function demonstrates how to:
    - Duplicate a slide and insert at a specific position
    - Add formatted code blocks to slides
    - Apply monospace font styling for code
    - Center-align title text

    CUSTOMIZATION:
    - Change source slide for duplication
    - Update target position (accounts for previously added slides)
    - Modify JSON code content
    - Change font (currently "Courier New")
    - Adjust font size (currently Pt(14))
    - Update title and subtitle text

    IMPORTANT: When adding multiple slides, remember that earlier additions
    shift the indices of later slides. This slide targets index 16 because
    the Ally McBeal slide was already inserted, shifting slide 16 to 17.

    Args:
        prs: pptx.presentation.Presentation object
    """
    print("  Duplicating slide 9 as template for JSON-LD slide...")

    # Duplicate slide 9 again as it has the right layout
    new_slide = duplicate_slide(prs, 8)
    print(f"  Created new slide at index {len(prs.slides) - 1}")

    # Move it to position 17 (index 16)
    print("  Moving to position 17...")
    slides = prs.slides._sldIdLst
    slide_element = slides[len(slides) - 1]
    slides.remove(slide_element)
    slides.insert(16, slide_element)

    # Update content in the moved slide
    print("  Updating slide content...")
    slide = prs.slides[16]

    # CUSTOMIZATION: Update JSON code content here
    # Multi-line string for code formatting
    json_code = '''{ "@type": "Product", "name": "Laptop", "offers": {
  "@type": "Offer", "price": "999.99", "priceCurrency": "GBP",
  "priceSpecification": [
    { "@type": "UnitPriceSpecification", "price": "899.99", "name": "Base Price" },

    { "@type": "DeliveryChargeSpecification", "price": "50.00", "name": "Delivery" },
    { "@type": "PaymentChargeSpecification", "price": "50.00", "name": "Commission" }
  ]
}'''

    # Identify shapes by their template content
    title_shape = None
    subtitle_shape = None
    body_shape = None

    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue

        text = shape.text_frame.text

        # Same identification strategy as Ally McBeal slide
        if "Mistake #1" in text:
            title_shape = shape
        elif "Pattern That Keeps" in text:
            subtitle_shape = shape
        elif "Toast notifications vanish" in text or "showToast" in text:
            body_shape = shape

    # Update title with center alignment
    if title_shape:
        print("    Updating title")
        title_shape.text_frame.clear()
        p = title_shape.text_frame.paragraphs[0]
        p.text = "Complete Pricing: JSON-LD Example"
        p.alignment = PP_ALIGN.CENTER  # Center the title
        # CUSTOMIZATION: Change alignment (LEFT, CENTER, RIGHT, JUSTIFY)

    # Update subtitle
    if subtitle_shape:
        print("    Updating subtitle")
        subtitle_shape.text_frame.clear()
        p = subtitle_shape.text_frame.paragraphs[0]
        p.text = "Machine-readable pricing with full cost breakdown, does not change human UI"

    # Update body with formatted code
    if body_shape:
        print("    Adding JSON code")
        body_shape.text_frame.clear()
        p = body_shape.text_frame.paragraphs[0]
        p.text = json_code

        # Apply monospace font for code appearance
        # CUSTOMIZATION: Change font and size for code blocks
        if p.runs:
            p.runs[0].font.name = "Courier New"  # Monospace font
            p.runs[0].font.size = Pt(14)  # 14-point size

    print("  ✓ JSON-LD slide added successfully")

if __name__ == "__main__":
    main()
